"""Rich document (PDF/Word/Excel/PowerPoint) to Markdown text extraction.

Each extractor converts one document format into Markdown-flavoured plain
text so the result can flow through the existing KnowledgeService chunking
pipeline (which splits on ``#`` headings and blank lines). Parser libraries
are imported lazily inside each extractor so environments without the
optional dependencies can still import this module.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import List

DOCUMENT_SUFFIXES = {".pdf", ".docx", ".xlsx", ".pptx"}
LEGACY_SUFFIXES = {".doc", ".xls", ".ppt"}

# Extraction budgets keep pathological documents from flooding the index.
MAX_SHEET_ROWS = 2000
MAX_TABLE_ROWS = 500
MAX_TOTAL_CHARACTERS = 400_000
TRUNCATION_NOTE = "\n\n（内容过长，已截断）"

_HEADING_STYLE = re.compile(r"^(?:heading|标题)\s*(\d)", re.IGNORECASE)


def extract_document_text(path: Path) -> str:
    """Convert a rich document into Markdown text, raising ``ValueError``
    with a user-facing Chinese message on any unsupported or broken input."""
    suffix = path.suffix.lower()
    if suffix in LEGACY_SUFFIXES:
        raise ValueError(
            "暂不支持旧版 Office 格式（.doc/.xls/.ppt），"
            "请先另存为 .docx/.xlsx/.pptx 再索引"
        )
    if suffix == ".pdf":
        text = _extract_pdf(path)
    elif suffix == ".docx":
        text = _extract_docx(path)
    elif suffix == ".xlsx":
        text = _extract_xlsx(path)
    elif suffix == ".pptx":
        text = _extract_pptx(path)
    else:
        raise ValueError("不支持的文档类型：{}".format(suffix or path.name))
    text = text.strip()
    if not text:
        raise ValueError("文档未包含可提取的文本内容")
    if len(text) > MAX_TOTAL_CHARACTERS:
        text = text[:MAX_TOTAL_CHARACTERS] + TRUNCATION_NOTE
    return text


def _cell_text(value: object) -> str:
    """Normalize one table cell for a Markdown table row."""
    if value is None:
        return ""
    text = str(value).replace("|", "\\|")
    return " ".join(text.split())


def _markdown_table(rows: List[List[str]]) -> str:
    """Render non-empty rows as a Markdown table with a header separator."""
    cleaned = [row for row in rows if any(cell for cell in row)]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    lines = []
    for index, row in enumerate(cleaned):
        padded = list(row) + [""] * (width - len(row))
        lines.append("| " + " | ".join(padded) + " |")
        if index == 0:
            lines.append("|" + " --- |" * width)
    return "\n".join(lines)


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - dependency is bundled
        raise ValueError("缺少 pypdf 依赖，无法解析 PDF 文件") from exc
    try:
        reader = PdfReader(str(path))
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception as exc:
                raise ValueError("PDF 已加密，无法提取文本") from exc
        pages: List[str] = []
        for number, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append("## 第{}页\n\n{}".format(number, text))
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError("无法解析该 PDF 文档：{}".format(exc)) from exc
    if not pages:
        raise ValueError("PDF 未包含可提取文本（可能是扫描件），请先转换为文字版")
    return "\n\n".join(pages)


def _docx_paragraph(paragraph) -> str:
    text = paragraph.text.strip()
    if not text:
        return ""
    style_name = getattr(getattr(paragraph, "style", None), "name", "") or ""
    match = _HEADING_STYLE.match(style_name)
    if match:
        level = min(max(int(match.group(1)), 1), 6)
        return "{} {}".format("#" * level, text)
    return text


def _docx_table(table) -> str:
    rows: List[List[str]] = []
    for row in table.rows[:MAX_TABLE_ROWS]:
        rows.append([_cell_text(cell.text) for cell in row.cells])
    return _markdown_table(rows)


def _extract_docx(path: Path) -> str:
    try:
        import docx
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:  # pragma: no cover - dependency is bundled
        raise ValueError("缺少 python-docx 依赖，无法解析 Word 文件") from exc
    try:
        document = docx.Document(str(path))
        blocks: List[str] = []
        # iter_inner_content keeps paragraphs and tables in document order.
        if hasattr(document, "iter_inner_content"):
            items: List[object] = list(document.iter_inner_content())
        else:  # pragma: no cover - older python-docx fallback
            items = list(document.paragraphs) + list(document.tables)
        for item in items:
            if isinstance(item, Paragraph):
                block = _docx_paragraph(item)
            elif isinstance(item, Table):
                block = _docx_table(item)
            else:  # pragma: no cover - future block types
                block = ""
            if block:
                blocks.append(block)
    except Exception as exc:
        raise ValueError("无法解析该 Word 文档：{}".format(exc)) from exc
    if not blocks:
        raise ValueError("Word 文档未包含可提取的文本内容")
    return "\n\n".join(blocks)


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover - dependency is bundled
        raise ValueError("缺少 openpyxl 依赖，无法解析 Excel 文件") from exc
    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
        sections: List[str] = []
        try:
            for sheet in workbook.worksheets:
                rows: List[List[str]] = []
                truncated = False
                for index, row in enumerate(sheet.iter_rows(values_only=True)):
                    if index >= MAX_SHEET_ROWS:
                        truncated = True
                        break
                    rows.append([_cell_text(value) for value in row])
                table = _markdown_table(rows)
                if not table:
                    continue
                section = "## {}\n\n{}".format(sheet.title, table)
                if truncated:
                    section += "\n\n（工作表行数过多，仅索引前 {} 行）".format(MAX_SHEET_ROWS)
                sections.append(section)
        finally:
            workbook.close()
    except Exception as exc:
        raise ValueError("无法解析该 Excel 文档：{}".format(exc)) from exc
    if not sections:
        raise ValueError("Excel 文档未包含可提取的内容")
    return "\n\n".join(sections)


def _pptx_table(shape) -> str:
    rows: List[List[str]] = []
    for row in shape.table.rows:
        rows.append([_cell_text(cell.text) for cell in row.cells])
    return _markdown_table(rows)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - dependency is bundled
        raise ValueError("缺少 python-pptx 依赖，无法解析 PPT 文件") from exc
    try:
        presentation = Presentation(str(path))
        slides: List[str] = []
        for number, slide in enumerate(presentation.slides, 1):
            parts: List[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
                if getattr(shape, "has_table", False):
                    table = _pptx_table(shape)
                    if table:
                        parts.append(table)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append("备注：{}".format(notes))
            if parts:
                slides.append("## 幻灯片{}\n\n{}".format(number, "\n\n".join(parts)))
    except Exception as exc:
        raise ValueError("无法解析该 PPT 文档：{}".format(exc)) from exc
    if not slides:
        raise ValueError("PPT 文档未包含可提取的文本内容")
    return "\n\n".join(slides)
