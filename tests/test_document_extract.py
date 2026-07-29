from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from src.core.services.document_extract import extract_document_text


def _minimal_pdf(text: str) -> bytes:
    """Build a one-page PDF by hand so tests do not need a PDF writer."""
    stream = "BT /F1 24 Tf 72 712 Td ({}) Tj ET".format(text).encode("latin-1")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n"
        + stream + b"\nendstream",
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = []
    for number, body in enumerate(objects, 1):
        offsets.append(len(output))
        output += "{} 0 obj\n".format(number).encode("ascii")
        output += body + b"\nendobj\n"
    xref_at = len(output)
    output += "xref\n0 {}\n".format(len(objects) + 1).encode("ascii")
    output += b"0000000000 65535 f \n"
    for offset in offsets:
        output += "{:010d} 00000 n \n".format(offset).encode("ascii")
    output += (
        "trailer\n<< /Size {} /Root 1 0 R >>\nstartxref\n{}\n%%EOF\n"
        .format(len(objects) + 1, xref_at).encode("ascii")
    )
    return bytes(output)


class DocumentExtractTests(unittest.TestCase):
    def setUp(self):
        self.temp = tempfile.TemporaryDirectory()
        self.addCleanup(self.temp.cleanup)
        self.root = Path(self.temp.name)

    def test_pdf_extracts_page_text_with_page_headings(self):
        path = self.root / "sample.pdf"
        path.write_bytes(_minimal_pdf("Hello Knowledge"))
        text = extract_document_text(path)
        self.assertIn("## 第1页", text)
        self.assertIn("Hello Knowledge", text)

    def test_pdf_without_text_raises(self):
        path = self.root / "blank.pdf"
        path.write_bytes(_minimal_pdf(""))
        with self.assertRaisesRegex(ValueError, "PDF 未包含可提取文本"):
            extract_document_text(path)

    def test_docx_headings_paragraphs_and_tables(self):
        import docx

        path = self.root / "manual.docx"
        document = docx.Document()
        document.add_heading("产品手册", level=1)
        document.add_paragraph("这是产品的详细说明。")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "功能"
        table.cell(0, 1).text = "说明"
        table.cell(1, 0).text = "检索"
        table.cell(1, 1).text = "支持混合检索"
        document.save(str(path))

        text = extract_document_text(path)
        self.assertIn("# 产品手册", text)
        self.assertIn("这是产品的详细说明。", text)
        self.assertIn("| 功能 | 说明 |", text)
        self.assertIn("| 检索 | 支持混合检索 |", text)

    def test_xlsx_sheets_become_markdown_tables(self):
        from openpyxl import Workbook

        path = self.root / "price.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "价格表"
        sheet.append(["商品", "价格"])
        sheet.append(["苹果", 5])
        workbook.save(str(path))

        text = extract_document_text(path)
        self.assertIn("## 价格表", text)
        self.assertIn("| 商品 | 价格 |", text)
        self.assertIn("| 苹果 | 5 |", text)

    def test_pptx_slides_text_and_notes(self):
        from pptx import Presentation

        path = self.root / "review.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "季度总结"
        slide.placeholders[1].text = "营收增长两成"
        slide.notes_slide.notes_text_frame.text = "记得补充图表"
        presentation.save(str(path))

        text = extract_document_text(path)
        self.assertIn("## 幻灯片1", text)
        self.assertIn("季度总结", text)
        self.assertIn("营收增长两成", text)
        self.assertIn("备注：记得补充图表", text)

    def test_legacy_office_format_is_rejected(self):
        path = self.root / "old.doc"
        path.write_bytes(b"\xd0\xcf\x11\xe0legacy")
        with self.assertRaisesRegex(ValueError, "旧版 Office"):
            extract_document_text(path)

    def test_unknown_suffix_is_rejected(self):
        path = self.root / "data.csv"
        path.write_text("a,b", encoding="utf-8")
        with self.assertRaisesRegex(ValueError, "不支持的文档类型"):
            extract_document_text(path)

    def test_broken_docx_raises_chinese_error(self):
        path = self.root / "broken.docx"
        path.write_bytes(b"not a real docx")
        with self.assertRaisesRegex(ValueError, "无法解析该 Word 文档"):
            extract_document_text(path)


if __name__ == "__main__":
    unittest.main()
